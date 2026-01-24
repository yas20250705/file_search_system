import os
import sqlite3
import time
import logging
import fitz  # PyMuPDF
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation
from datetime import datetime

from database import get_index_db_connection, update_indexing_status, is_indexing_stop_requested, set_indexing_stop_requested, index_db_lock, update_index_status

logger = logging.getLogger(__name__)

# --- Text Extraction Functions ---

def extract_text_from_pdf(file_path):
    logger.debug(f"PDF抽出開始: {file_path}")
    try:
        with fitz.open(file_path) as doc:
            logger.debug(f"PDFファイルオープン成功: {file_path}")
            text = "".join(page.get_text() for page in doc)
        logger.debug(f"PDF抽出完了: {file_path}")
        return text
    except Exception as e:
        logger.error(f"PDFファイルからのテキスト抽出エラー ({file_path}): {e}", exc_info=True)
        return ""

def extract_text_from_excel(file_path):
    try:
        workbook = load_workbook(filename=file_path, read_only=True)
        text = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        text.append(str(cell.value))
        return " ".join(text)
    except Exception as e:
        logger.error(f"Excelファイルからのテキスト抽出エラー ({file_path}): {e}")
        return ""

def extract_text_from_word(file_path):
    try:
        doc = Document(file_path)
        text = [p.text for p in doc.paragraphs]
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Wordファイルからのテキスト抽出エラー ({file_path}): {e}")
        return ""

def extract_text_from_powerpoint(file_path):
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): # テキストを持つシェイプのみ
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logger.error(f"PowerPointファイルからのテキスト抽出エラー ({file_path}): {e}")
        return ""

def extract_text_from_plain(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        logger.error(f"テキストファイルからの読み込みエラー ({file_path}): {e}")
        return ""

# --- Main Indexing Logic ---

def index_files(index_id: int, target_directory: str, allowed_extensions: list[str], db_path: str):
    logger.info(f"インデックスID {index_id} ('{target_directory}') のインデックス作成を開始します...")
    start_time = time.time()
    
    # メタデータベースのステータスを更新
    update_index_status(index_id, 'running')

    conn = None # 接続を初期化
    try:
        conn = get_index_db_connection(db_path)
        cursor = conn.cursor()

        # インデックス作成時は常にテーブルを完全に削除して再作成
        # これにより、FTS5のcontent-syncテーブルの同期問題やトークナイザーの問題を回避
        logger.info(f"インデックスID {index_id} のテーブルを再作成します...")
        
        # テーブルを完全に削除
        cursor.execute("DROP TABLE IF EXISTS files_fts")
        cursor.execute("DROP TABLE IF EXISTS files")
        conn.commit()
        
        # filesテーブルを再作成
        cursor.execute("""
            CREATE TABLE files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                path TEXT NOT NULL UNIQUE,
                content TEXT,
                file_type TEXT,
                modified_date REAL,
                created_date REAL
            )
        """)
        # FTS5テーブルをtrigramトークナイザーで作成（content-syncを使用しない）
        # content-syncとtrigramの組み合わせは問題を引き起こすため、独立したテーブルを使用
        cursor.execute("""
            CREATE VIRTUAL TABLE files_fts USING fts5(
                path, 
                content,
                tokenize = 'trigram'
            )
        """)
        conn.commit()
        logger.info(f"インデックスID {index_id} のテーブルをtrigramトークナイザーで再作成しました。")

        files_to_index = []
        for root, _, files in os.walk(target_directory):
            for file in files:
                logger.debug(f"Indexer: Found file: {os.path.join(root, file)}")
                if any(file.endswith(ext) for ext in allowed_extensions):
                    files_to_index.append(os.path.join(root, file))
        
        logger.debug(f"Indexer: Files to index after filtering: {files_to_index}")
        total_files = len(files_to_index)
        logger.info(f"インデックスID {index_id} の対象ファイル数: {total_files}")
        
        update_indexing_status(conn, db_path, "started", total_files, 0, start_time, 0) # 個別DBのステータスを更新

        if total_files == 0:
            logger.info(f"インデックスID {index_id} の対象ファイルがありません。インデックス作成を完了します。")
            with index_db_lock:
                update_indexing_status(conn, db_path, "completed", 0, 0, start_time, time.time()) # 個別DBのステータスを更新
            update_index_status(index_id, 'completed', datetime.now())
            return # 関数を終了

        logger.debug(f"Indexer: Starting file processing loop for {total_files} files.")
        for i, file_path in enumerate(files_to_index):
            if is_indexing_stop_requested(conn, db_path):
                logger.info(f"インデックスID {index_id} のインデックス作成がユーザーによって中止されました。")
                update_indexing_status(conn, db_path, "stopped", total_files, i, start_time, time.time()) # 個別DBのステータスを更新
                update_index_status(index_id, 'stopped') # メタDBのステータスを更新
                break

            ext = os.path.splitext(file_path)[1].lower()
            content = ""
            logger.debug(f"Indexer: Extracting text from {file_path}")
            if ext == '.pdf':
                content = extract_text_from_pdf(file_path)
            elif ext in ['.xlsx', '.xls']:
                content = extract_text_from_excel(file_path)
            elif ext == '.docx':
                content = extract_text_from_word(file_path)
            elif ext == '.pptx':
                content = extract_text_from_powerpoint(file_path)
            else:
                content = extract_text_from_plain(file_path)
            logger.debug(f"Indexer: Finished extracting text from {file_path}")

            # ファイル情報を取得
            file_type = ext
            # 指定したフォルダ内のファイルの日時を使用（ファイルシステムの日時）
            try:
                modified_timestamp = os.path.getmtime(file_path)
                created_timestamp = os.path.getctime(file_path)
            except OSError as e:
                logger.warning(f"ファイル情報の取得に失敗しました ({file_path}): {e}")
                modified_timestamp = None
                created_timestamp = None

            # contentが空でもファイル情報は保存する
            try:
                # 1. `files`テーブルに挿入（ファイル情報を含む）
                # contentが空の場合は空文字列を保存
                content_to_save = content if content else ""
                
                # 新規レコードを挿入（テーブルは毎回再作成されるので既存レコードはない）
                cursor.execute("INSERT INTO files (path, content, file_type, modified_date, created_date) VALUES (?, ?, ?, ?, ?)", 
                             (file_path, content_to_save, file_type, modified_timestamp, created_timestamp))
                
                # 2. `files_fts`テーブルに挿入（content-syncを使用しない独立したテーブル）
                if content:
                    cursor.execute("INSERT INTO files_fts (path, content) VALUES (?, ?)", (file_path, content))
                        
            except sqlite3.Error as e:
                logger.error(f"インデックスID {index_id} のデータベース挿入エラー ({file_path}): {e}")

            # 進捗を更新
            current_processed_files = i + 1
            logger.debug(f"Indexer: Calling update_indexing_status for index {index_id} with processed_files={current_processed_files}/{total_files}")
            update_indexing_status(conn, db_path, "running", total_files, current_processed_files, start_time, 0) # 個別DBのステータスを更新

            if current_processed_files % 10 == 0:
                conn.commit() # 10ファイルごとにコミット
                logger.info(f"インデックスID {index_id} の進捗: {current_processed_files}/{total_files}")

        conn.commit() # 最終コミット
        
        if not is_indexing_stop_requested(conn, db_path): # 中止されていない場合のみ完了ステータス
            logger.info(f"インデックスID {index_id} のインデックス作成が完了しました。")
            update_indexing_status(conn, db_path, "completed", total_files, total_files, start_time, time.time()) # 個別DBのステータスを更新
            update_index_status(index_id, 'completed', datetime.now()) # メタDBのステータスと最終インデックス日時を更新
        else:
            logger.info(f"インデックスID {index_id} のインデックス作成は中止されました。")
            update_indexing_status(conn, db_path, "stopped", total_files, i, start_time, time.time()) # 個別DBのステータスを更新
            update_index_status(index_id, 'stopped') # メタDBのステータスを更新

    except Exception as e:
        logger.error(f"インデックスID {index_id} のインデックス作成中に予期せぬエラーが発生しました: {e}", exc_info=True)
        update_indexing_status(conn, db_path, "failed", total_files, i, start_time, time.time()) # 個別DBのステータスを更新
        update_index_status(index_id, 'failed') # メタDBのステータスを更新
    finally:
        if conn:
            conn.close()


def extract_content(file_path: str) -> str:
    """ファイルからテキストを抽出するヘルパー関数"""
    ext = os.path.splitext(file_path)[1].lower()
    content = ""
    
    if ext == '.pdf':
        content = extract_text_from_pdf(file_path)
    elif ext in ['.xlsx', '.xls']:
        content = extract_text_from_excel(file_path)
    elif ext == '.docx':
        content = extract_text_from_word(file_path)
    elif ext == '.pptx':
        content = extract_text_from_powerpoint(file_path)
    else:
        content = extract_text_from_plain(file_path)
    
    return content


def update_index_files(index_id: int, target_directory: str, allowed_extensions: list[str], db_path: str):
    """
    インデックスを差分更新します。
    - 新規ファイル: INSERT
    - 更新ファイル（タイムスタンプが異なる）: UPDATE
    - 削除ファイル（ディスクにない）: DELETE
    - 変更なし: スキップ
    """
    logger.info(f"インデックスID {index_id} ('{target_directory}') の差分更新を開始します...")
    start_time = time.time()
    
    # メタデータベースのステータスを更新
    update_index_status(index_id, 'running')
    
    conn = None
    processed_count = 0
    total_files = 0
    
    try:
        conn = get_index_db_connection(db_path)
        cursor = conn.cursor()
        
        # テーブルが存在するか確認
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='files'")
        if not cursor.fetchone():
            logger.warning(f"インデックスID {index_id} のfilesテーブルが存在しません。完全インデックスを実行してください。")
            update_index_status(index_id, 'failed')
            return
        
        # 1. DBから既存ファイル情報を取得
        cursor.execute("SELECT path, modified_date FROM files")
        existing_files = {row[0]: row[1] for row in cursor.fetchall()}
        logger.info(f"インデックスID {index_id} の既存ファイル数: {len(existing_files)}")
        
        # 2. ディレクトリをスキャンして現在のファイルリストを取得
        current_files = []
        for root, _, files in os.walk(target_directory):
            for file in files:
                if any(file.endswith(ext) for ext in allowed_extensions):
                    current_files.append(os.path.join(root, file))
        
        current_files_set = set(current_files)
        existing_files_set = set(existing_files.keys())
        
        # 3. 差分を検出
        new_files = current_files_set - existing_files_set
        deleted_files = existing_files_set - current_files_set
        potentially_updated_files = current_files_set & existing_files_set
        
        # 更新されたファイルを検出（タイムスタンプが異なるもの）
        updated_files = []
        for file_path in potentially_updated_files:
            try:
                current_mtime = os.path.getmtime(file_path)
                stored_mtime = existing_files.get(file_path)
                # タイムスタンプが異なる場合は更新対象
                if stored_mtime is None or abs(current_mtime - stored_mtime) > 1:  # 1秒の誤差を許容
                    updated_files.append(file_path)
            except OSError:
                pass  # ファイルにアクセスできない場合はスキップ
        
        total_files = len(new_files) + len(updated_files) + len(deleted_files)
        logger.info(f"インデックスID {index_id} の差分: 新規={len(new_files)}, 更新={len(updated_files)}, 削除={len(deleted_files)}, 変更なし={len(potentially_updated_files) - len(updated_files)}")
        
        if total_files == 0:
            logger.info(f"インデックスID {index_id} の更新対象ファイルがありません。")
            update_indexing_status(conn, db_path, "completed", 0, 0, start_time, time.time())
            update_index_status(index_id, 'completed', datetime.now())
            return
        
        update_indexing_status(conn, db_path, "started", total_files, 0, start_time, 0)
        
        # 4. 削除ファイルを処理
        for file_path in deleted_files:
            if is_indexing_stop_requested(conn, db_path):
                logger.info(f"インデックスID {index_id} の更新がユーザーによって中止されました。")
                update_indexing_status(conn, db_path, "stopped", total_files, processed_count, start_time, time.time())
                update_index_status(index_id, 'stopped')
                return
            
            try:
                cursor.execute("DELETE FROM files WHERE path = ?", (file_path,))
                cursor.execute("DELETE FROM files_fts WHERE path = ?", (file_path,))
                logger.debug(f"削除: {file_path}")
            except sqlite3.Error as e:
                logger.error(f"削除エラー ({file_path}): {e}")
            
            processed_count += 1
            update_indexing_status(conn, db_path, "running", total_files, processed_count, start_time, 0)
        
        # 5. 新規ファイルを処理
        for file_path in new_files:
            if is_indexing_stop_requested(conn, db_path):
                logger.info(f"インデックスID {index_id} の更新がユーザーによって中止されました。")
                update_indexing_status(conn, db_path, "stopped", total_files, processed_count, start_time, time.time())
                update_index_status(index_id, 'stopped')
                return
            
            try:
                ext = os.path.splitext(file_path)[1].lower()
                content = extract_content(file_path)
                
                modified_timestamp = os.path.getmtime(file_path)
                created_timestamp = os.path.getctime(file_path)
                
                content_to_save = content if content else ""
                
                cursor.execute(
                    "INSERT INTO files (path, content, file_type, modified_date, created_date) VALUES (?, ?, ?, ?, ?)",
                    (file_path, content_to_save, ext, modified_timestamp, created_timestamp)
                )
                
                if content:
                    cursor.execute("INSERT INTO files_fts (path, content) VALUES (?, ?)", (file_path, content))
                
                logger.debug(f"新規追加: {file_path}")
            except sqlite3.Error as e:
                logger.error(f"新規追加エラー ({file_path}): {e}")
            except OSError as e:
                logger.warning(f"ファイル情報取得エラー ({file_path}): {e}")
            
            processed_count += 1
            update_indexing_status(conn, db_path, "running", total_files, processed_count, start_time, 0)
            
            if processed_count % 10 == 0:
                conn.commit()
                logger.info(f"インデックスID {index_id} の進捗: {processed_count}/{total_files}")
        
        # 6. 更新ファイルを処理
        for file_path in updated_files:
            if is_indexing_stop_requested(conn, db_path):
                logger.info(f"インデックスID {index_id} の更新がユーザーによって中止されました。")
                update_indexing_status(conn, db_path, "stopped", total_files, processed_count, start_time, time.time())
                update_index_status(index_id, 'stopped')
                return
            
            try:
                ext = os.path.splitext(file_path)[1].lower()
                content = extract_content(file_path)
                
                modified_timestamp = os.path.getmtime(file_path)
                created_timestamp = os.path.getctime(file_path)
                
                content_to_save = content if content else ""
                
                # filesテーブルを更新
                cursor.execute(
                    "UPDATE files SET content = ?, file_type = ?, modified_date = ?, created_date = ? WHERE path = ?",
                    (content_to_save, ext, modified_timestamp, created_timestamp, file_path)
                )
                
                # files_ftsテーブルを更新（一度削除して再挿入）
                cursor.execute("DELETE FROM files_fts WHERE path = ?", (file_path,))
                if content:
                    cursor.execute("INSERT INTO files_fts (path, content) VALUES (?, ?)", (file_path, content))
                
                logger.debug(f"更新: {file_path}")
            except sqlite3.Error as e:
                logger.error(f"更新エラー ({file_path}): {e}")
            except OSError as e:
                logger.warning(f"ファイル情報取得エラー ({file_path}): {e}")
            
            processed_count += 1
            update_indexing_status(conn, db_path, "running", total_files, processed_count, start_time, 0)
            
            if processed_count % 10 == 0:
                conn.commit()
                logger.info(f"インデックスID {index_id} の進捗: {processed_count}/{total_files}")
        
        conn.commit()
        
        if not is_indexing_stop_requested(conn, db_path):
            logger.info(f"インデックスID {index_id} の差分更新が完了しました。")
            update_indexing_status(conn, db_path, "completed", total_files, total_files, start_time, time.time())
            update_index_status(index_id, 'completed', datetime.now())
        else:
            logger.info(f"インデックスID {index_id} の差分更新は中止されました。")
            update_indexing_status(conn, db_path, "stopped", total_files, processed_count, start_time, time.time())
            update_index_status(index_id, 'stopped')
    
    except Exception as e:
        logger.error(f"インデックスID {index_id} の差分更新中に予期せぬエラーが発生しました: {e}", exc_info=True)
        update_indexing_status(conn, db_path, "failed", total_files, processed_count, start_time, time.time())
        update_index_status(index_id, 'failed')
    finally:
        if conn:
            conn.close()